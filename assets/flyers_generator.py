import tkinter as tk
from tkinter import ttk, scrolledtext, filedialog
import threading
import os
import re
import base64
import pythoncom
import win32com.client
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from deep_translator import GoogleTranslator
from PIL import Image

# Dimensions cibles pour les visuels de flyer, imposees par le rectangle QR
# deja calibre dans redirect-template.js (window.FLYER_TEMPLATES[...].rect).
# Le rectangle QR est fixe en pixels absolus (x:27, y:154, width:292, height:293)
# sur un visuel de 1600x896 : si on change cette taille ici, il faut aussi
# recalculer/mettre a jour "rect" pour chaque langue dans redirect-template.js.
FLYER_TARGET_SIZE = (1600, 896)

# Chemin par defaut vers redirect-template.js (a la racine du repo, un niveau
# au-dessus de ce script qui vit dans assets/). Modifiable dans l'interface.
DEFAULT_TEMPLATE_JS_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "redirect-template.js")

# Libelles affiches dans le selecteur de flyer (url-generator.html) pour
# chaque code langue. Complete cette liste si tu ajoutes une langue via le
# champ "Autres langues" — sinon le code (ex: "PT") sera utilise tel quel.
LANG_LABELS = {
    "fr": "Français", "en": "English", "de": "Deutsch", "nl": "Nederlands",
    "es": "Español", "it": "Italiano", "da": "Dansk", "sv": "Svenska",
    "pt": "Português", "ru": "Русский", "pl": "Polski",
}

class PPTXTranslatorApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Traducteur & Extracteur PPTX - Bell'O du Lac")
        self.geometry("700x820")
        self.configure(padx=20, pady=10)

        # ==========================================
        # SECTION 1 : CONFIGURATION (Fichier & Langues)
        # ==========================================
        frame_src = tk.LabelFrame(self, text=" ⚙️ Configuration ", font=("Arial", 11, "bold"), padx=15, pady=10)
        frame_src.pack(fill="x", pady=5)

        tk.Label(frame_src, text="Fichier source (FR) :").pack(anchor="w")
        self.source_entry = tk.Entry(frame_src, width=50)
        self.source_entry.insert(0, "flyers.pptx")
        self.source_entry.pack(fill="x", pady=(0, 15))

        tk.Label(frame_src, text="Langues à générer :").pack(anchor="w")
        
        self.lang_vars = {
            "en": tk.BooleanVar(value=True),
            "de": tk.BooleanVar(value=True),
            "nl": tk.BooleanVar(value=True),
            "es": tk.BooleanVar(value=True),
            "it": tk.BooleanVar(value=True),
            "da": tk.BooleanVar(value=True),
            "sv": tk.BooleanVar(value=True)
        }
        
        lang_frame1 = tk.Frame(frame_src)
        lang_frame1.pack(fill="x")
        tk.Checkbutton(lang_frame1, text="Anglais (EN)", variable=self.lang_vars["en"]).pack(side="left", padx=5)
        tk.Checkbutton(lang_frame1, text="Allemand (DE)", variable=self.lang_vars["de"]).pack(side="left", padx=5)
        tk.Checkbutton(lang_frame1, text="Néerlandais (NL)", variable=self.lang_vars["nl"]).pack(side="left", padx=5)
        tk.Checkbutton(lang_frame1, text="Espagnol (ES)", variable=self.lang_vars["es"]).pack(side="left", padx=5)

        lang_frame2 = tk.Frame(frame_src)
        lang_frame2.pack(fill="x", pady=(5, 10))
        tk.Checkbutton(lang_frame2, text="Italien (IT)", variable=self.lang_vars["it"]).pack(side="left", padx=5)
        tk.Checkbutton(lang_frame2, text="Danois (DA)", variable=self.lang_vars["da"]).pack(side="left", padx=5)
        tk.Checkbutton(lang_frame2, text="Suédois (SV)", variable=self.lang_vars["sv"]).pack(side="left", padx=5)

        tk.Label(frame_src, text="Autres langues (ex: pt, ru, pl) :").pack(anchor="w")
        self.other_langs_entry = tk.Entry(frame_src, width=50)
        self.other_langs_entry.pack(fill="x", pady=(0, 15))

        tk.Label(frame_src, text="Fichier redirect-template.js à mettre à jour automatiquement :").pack(anchor="w")
        template_js_frame = tk.Frame(frame_src)
        template_js_frame.pack(fill="x", pady=(0, 5))
        self.template_js_entry = tk.Entry(template_js_frame)
        self.template_js_entry.insert(0, DEFAULT_TEMPLATE_JS_PATH)
        self.template_js_entry.pack(side="left", fill="x", expand=True, padx=(0, 10))
        tk.Button(template_js_frame, text="Parcourir...", command=self.browse_template_js).pack(side="right")
        tk.Label(frame_src, text="Laisser vide pour ne pas toucher à redirect-template.js (juste générer les JPG/base64).", font=("Arial", 8, "italic"), fg="#555").pack(anchor="w", pady=(0, 5))

        # ==========================================
        # SECTION 2 : ACTIONS
        # ==========================================
        frame_actions = tk.LabelFrame(self, text=" ⚡ Actions ", font=("Arial", 11, "bold"), padx=15, pady=10)
        frame_actions.pack(fill="x", pady=5)

        self.btn_all = tk.Button(frame_actions, text="🚀 TOUT GÉNÉRER D'UN COUP (Traduction + Images + Base64)", 
                                 bg="#9C27B0", fg="white", font=("Arial", 10, "bold"), command=self.start_generate_all)
        self.btn_all.pack(fill="x", pady=(0, 15))
        
        ttk.Separator(frame_actions, orient='horizontal').pack(fill='x', pady=5)
        tk.Label(frame_actions, text="Ou en deux étapes (si vous souhaitez corriger le texte entre temps) :", font=("Arial", 9, "italic")).pack(anchor="w", pady=(5, 5))

        self.btn_translate = tk.Button(frame_actions, text="1. Générer uniquement les PPTX", 
                                       bg="#2196F3", fg="white", font=("Arial", 9, "bold"), command=self.start_translation)
        self.btn_translate.pack(fill="x", pady=5)

        export_frame = tk.Frame(frame_actions)
        export_frame.pack(fill="x", pady=5)
        self.export_entry = tk.Entry(export_frame)
        self.export_entry.pack(side="left", fill="x", expand=True, padx=(0, 10))
        tk.Button(export_frame, text="Parcourir...", command=self.browse_file).pack(side="right")
        
        self.btn_export = tk.Button(frame_actions, text="2. Exporter Images & Base64 du PPTX sélectionné", 
                                    bg="#4CAF50", fg="white", font=("Arial", 9, "bold"), command=self.start_export)
        self.btn_export.pack(fill="x", pady=5)

        # ==========================================
        # SECTION 3 : CONSOLE
        # ==========================================
        tk.Label(self, text="Journal d'activité :", font=("Arial", 10, "bold")).pack(anchor="w")
        self.console = scrolledtext.ScrolledText(self, height=10, bg="#f4f4f4")
        self.console.pack(fill="both", expand=True)

    def log(self, message):
        self.console.insert(tk.END, message + "\n")
        self.console.see(tk.END)
        self.update_idletasks()

    def browse_file(self):
        filename = filedialog.askopenfilename(filetypes=[("PowerPoint", "*.pptx")])
        if filename:
            self.export_entry.delete(0, tk.END)
            self.export_entry.insert(0, filename)

    def browse_template_js(self):
        filename = filedialog.askopenfilename(filetypes=[("JavaScript", "*.js")])
        if filename:
            self.template_js_entry.delete(0, tk.END)
            self.template_js_entry.insert(0, filename)

    def set_buttons_state(self, state):
        self.btn_all.config(state=state)
        self.btn_translate.config(state=state)
        self.btn_export.config(state=state)

    def get_selected_langs(self):
        langs = [code for code, var in self.lang_vars.items() if var.get()]
        other_langs = self.other_langs_entry.get().strip()
        if other_langs:
            langs.extend([lang.strip() for lang in other_langs.split(",") if lang.strip()])
        return [l.lower() for l in langs]

    # --- NOYAU DE TRADUCTION AVEC GESTION DES ESPACES ---
    def translate_text(self, text, target_lang):
        if not text or not text.strip(): return text
        if not any(c.isalpha() for c in text): return text
        
        # 1. Capture mathématique des espaces de début et de fin
        len_leading = len(text) - len(text.lstrip())
        len_trailing = len(text) - len(text.rstrip())
        
        prefix = text[:len_leading] if len_leading > 0 else ""
        suffix = text[-len_trailing:] if len_trailing > 0 else ""
        
        # 2. Le texte "pur" à envoyer au traducteur
        core_text = text.strip()
        
        try:
            translated = GoogleTranslator(source='fr', target=target_lang).translate(core_text)
            if translated:
                # 3. On recolle les espaces initiaux autour de la traduction
                return prefix + translated + suffix
            return text
        except Exception as e:
            self.log(f"  [!] Erreur traduction sur '{text}': {e}")
            return text

    def explore_and_translate_shape(self, shape, target_lang):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                self.explore_and_translate_shape(sub_shape, target_lang)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            translated = self.translate_text(run.text, target_lang)
                            if translated: run.text = translated
        elif shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    translated = self.translate_text(run.text, target_lang)
                    if translated: run.text = translated

    def perform_translation(self, source_file, target_lang):
        self.log(f"\n🔄 Traduction vers {target_lang.upper()} en cours...")
        try:
            prs = Presentation(source_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    self.explore_and_translate_shape(shape, target_lang)
            
            base_name, _ = os.path.splitext(source_file)
            output_pptx = f"{base_name}_{target_lang.upper()}.pptx"
            prs.save(output_pptx)
            self.log(f"✅ Fichier créé : {output_pptx}")
            return output_pptx
        except Exception as e:
            self.log(f"❌ Erreur lors de la traduction : {e}")
            return None

    # --- NOYAU D'EXPORT ---
    def resize_flyer_to_target(self, jpg_path):
        """Recadre l'export PowerPoint (taille variable selon la resolution
        d'ecran/DPI) a la taille fixe attendue par le rectangle QR cote site
        (voir FLYER_TARGET_SIZE et redirect-template.js). A appeler juste
        apres shape_range.Export(), avant la generation du base64."""
        try:
            with Image.open(jpg_path) as img:
                if img.size != FLYER_TARGET_SIZE:
                    resized = img.convert("RGB").resize(FLYER_TARGET_SIZE, Image.LANCZOS)
                    resized.save(jpg_path, quality=90)
                    self.log(f"  🔧 Redimensionné {img.size} → {FLYER_TARGET_SIZE[0]}x{FLYER_TARGET_SIZE[1]} (pour coller au QR code)")
        except Exception as e:
            self.log(f"  [!] Erreur redimensionnement : {e}")

    def generate_base64_from_jpg(self, jpg_path):
        """Ecrit le fichier _base64.txt (comme avant) et retourne le data URI
        genere, pour que perform_export puisse l'injecter dans redirect-template.js."""
        txt_path = jpg_path.replace(".jpg", "_base64.txt")
        try:
            with open(jpg_path, "rb") as img_file:
                b64_string = base64.b64encode(img_file.read()).decode('utf-8')
            data_uri = f"data:image/jpeg;base64,{b64_string}"
            with open(txt_path, "w", encoding="utf-8") as txt_file:
                txt_file.write(data_uri)
            return data_uri
        except Exception as e:
            self.log(f"  [!] Erreur Base64 : {e}")
            return None

    def infer_lang_code(self, pptx_path):
        """Devine le code langue a partir du nom de fichier genere par
        perform_translation (ex: 'flyers_DE.pptx' -> 'de'). Si le fichier n'a
        pas de suffixe langue (le pptx source, ex: 'flyers.pptx'), on suppose
        que c'est la version FR (libellee comme telle dans la Configuration)."""
        base_name, _ = os.path.splitext(os.path.basename(pptx_path))
        m = re.search(r'_([A-Za-z]{2,3})$', base_name)
        if m:
            return m.group(1).lower()
        return "fr"

    # --- MISE A JOUR AUTOMATIQUE DE redirect-template.js ---
    def update_flyer_template_js(self, js_path, lang_code, data_uri):
        """Insere ou remplace l'entree `lang_code` dans window.FLYER_TEMPLATES
        au sein de redirect-template.js, avec le nouveau visuel en base64.
        - Si la langue existe deja : on ne touche qu'a son "image" (le label
          et le rect existants sont conserves).
        - Si c'est une nouvelle langue : on ajoute une entree juste avant la
          fermeture de FLYER_TEMPLATES, en reprenant le meme "rect" que les
          autres langues (toutes calibrees sur le meme visuel 1600x896)."""
        if not js_path:
            return
        if not os.path.exists(js_path):
            self.log(f"  [!] redirect-template.js introuvable : {js_path}")
            return

        try:
            with open(js_path, "r", encoding="utf-8") as f:
                content = f.read()
        except Exception as e:
            self.log(f"  [!] Impossible de lire redirect-template.js : {e}")
            return

        # 1) La langue existe deja : on remplace juste son image, on garde
        #    label/rect intacts.
        pattern_existing = re.compile(r'(\n  ' + re.escape(lang_code) + r': \{\n    image: ")[^"]*(")')
        new_content, n = pattern_existing.subn(lambda m: m.group(1) + data_uri + m.group(2), content, count=1)

        if n > 0:
            with open(js_path, "w", encoding="utf-8") as f:
                f.write(new_content)
            self.log(f"  🔗 redirect-template.js : image '{lang_code}' mise à jour.")
            return

        # 2) Nouvelle langue : on localise le bloc "window.FLYER_TEMPLATES = { ... };"
        #    par comptage d'accolades (plus fiable qu'une regex sur tout le
        #    fichier, qui pourrait contenir d'autres objets ailleurs).
        start_match = re.search(r'window\.FLYER_TEMPLATES\s*=\s*\{', content)
        if not start_match:
            self.log("  [!] Bloc window.FLYER_TEMPLATES introuvable dans redirect-template.js.")
            return

        brace_start = start_match.end() - 1  # position du '{' d'ouverture
        depth = 0
        end_idx = None
        for i in range(brace_start, len(content)):
            if content[i] == '{':
                depth += 1
            elif content[i] == '}':
                depth -= 1
                if depth == 0:
                    end_idx = i
                    break
        if end_idx is None:
            self.log("  [!] Accolade fermante de FLYER_TEMPLATES introuvable.")
            return

        # Reprend le rect de la premiere langue existante (toutes identiques).
        rect_match = re.search(r'rect: \{[^}]*\}', content[brace_start:end_idx])
        rect_str = rect_match.group(0) if rect_match else 'rect: { x: 27, y: 154, width: 292, height: 293 }'
        label = LANG_LABELS.get(lang_code, lang_code.upper())

        entry = (
            f'\n  {lang_code}: {{\n'
            f'    image: "{data_uri}",\n'
            f'    label: "{label}",\n'
            f'    {rect_str}\n'
            f'  }},'
        )

        # La derniere entree existante n'a pas forcement de virgule finale
        # (dernier element d'un objet litteral) : on en ajoute une avant
        # d'inserer la nouvelle entree, sinon le JS devient invalide.
        before = content[:end_idx].rstrip()
        separator = '' if before.endswith(',') else ','
        new_content = before + separator + entry + "\n" + content[end_idx:]
        with open(js_path, "w", encoding="utf-8") as f:
            f.write(new_content)
        self.log(f"  🔗 redirect-template.js : nouvelle langue '{lang_code}' ajoutée (label \"{label}\").")

    def perform_export(self, ppt_app, pptx_path):
        self.log(f"📸 Export des objets depuis : {os.path.basename(pptx_path)}")
        abs_pptx_path = os.path.abspath(pptx_path)
        abs_dir = os.path.dirname(abs_pptx_path)
        base_name, _ = os.path.splitext(os.path.basename(pptx_path))
        lang_code = self.infer_lang_code(pptx_path)
        template_js_path = self.template_js_entry.get().strip()

        try:
            presentation = ppt_app.Presentations.Open(abs_pptx_path, WithWindow=False)
            for slide_idx, slide in enumerate(presentation.Slides, start=1):
                jpg_filename = f"{base_name}_slide{slide_idx}.jpg"
                abs_jpg_path = os.path.join(abs_dir, jpg_filename)
                try:
                    shape_range = slide.Shapes.Range()
                    shape_range.Export(abs_jpg_path, 1) # 1 = Filtre JPG
                    self.log(f"  ✅ Image exportée : {jpg_filename}")
                    self.resize_flyer_to_target(abs_jpg_path)
                    data_uri = self.generate_base64_from_jpg(abs_jpg_path)
                    if data_uri and template_js_path:
                        self.update_flyer_template_js(template_js_path, lang_code, data_uri)
                except Exception as e:
                    self.log(f"  [!] Erreur sur le slide {slide_idx} : {e}")
            presentation.Close()
        except Exception as e:
            self.log(f"❌ Erreur d'ouverture dans PowerPoint : {e}")

    # --- THREADS D'EXECUTION ---
    def start_generate_all(self):
        self.set_buttons_state("disabled")
        self.console.delete(1.0, tk.END)
        threading.Thread(target=self.generate_all_task, daemon=True).start()

    def generate_all_task(self):
        pythoncom.CoInitialize()
        ppt_app = None
        try:
            source = self.source_entry.get().strip()
            if not os.path.exists(source):
                self.log(f"❌ Fichier source '{source}' introuvable.")
                return

            langs = self.get_selected_langs()
            if not langs:
                self.log("❌ Veuillez sélectionner au moins une langue.")
                return

            self.log("🔧 Démarrage de PowerPoint en arrière-plan...")
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")

            for lang in langs:
                self.log(f"\n=============================================")
                self.log(f"🚀 TRAITEMENT COMPLET : {lang.upper()}")
                self.log(f"=============================================")
                
                # 1. Traduction
                out_pptx = self.perform_translation(source, lang)
                
                # 2. Export si traduction réussie
                if out_pptx:
                    self.perform_export(ppt_app, out_pptx)

            self.log("\n🎉 TOUT EST TERMINÉ ! Images et Base64 générés avec succès. 🎉")

        except Exception as e:
            self.log(f"\n❌ Erreur inattendue : {e}")
        finally:
            if ppt_app:
                try: ppt_app.Quit()
                except: pass
            pythoncom.CoUninitialize()
            self.set_buttons_state("normal")

    def start_translation(self):
        self.set_buttons_state("disabled")
        self.console.delete(1.0, tk.END)
        threading.Thread(target=self.translation_only_task, daemon=True).start()

    def translation_only_task(self):
        try:
            source = self.source_entry.get().strip()
            if not os.path.exists(source):
                self.log(f"❌ Fichier '{source}' introuvable.")
                return

            langs = self.get_selected_langs()
            if not langs: return

            for lang in langs:
                self.perform_translation(source, lang)
                
            self.log("\n🎉 Traduction terminée ! Éditez les PPTX avant l'étape 2.")
        except Exception as e:
            self.log(f"\n❌ Erreur : {e}")
        finally:
            self.set_buttons_state("normal")

    def start_export(self):
        self.set_buttons_state("disabled")
        self.console.delete(1.0, tk.END)
        threading.Thread(target=self.export_only_task, daemon=True).start()

    def export_only_task(self):
        pythoncom.CoInitialize()
        ppt_app = None
        try:
            filepath = self.export_entry.get().strip()
            if not os.path.exists(filepath):
                self.log(f"❌ Fichier à exporter '{filepath}' introuvable.")
                return

            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            self.perform_export(ppt_app, filepath)
            self.log("\n🎉 Export terminé avec succès !")
        except Exception as e:
            self.log(f"\n❌ Erreur : {e}")
        finally:
            if ppt_app:
                try: ppt_app.Quit()
                except: pass
            pythoncom.CoUninitialize()
            self.set_buttons_state("normal")

if __name__ == "__main__":
    app = PPTXTranslatorApp()
    app.mainloop()